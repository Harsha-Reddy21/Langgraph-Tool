import os
import requests
from bs4 import BeautifulSoup
from typing import Dict, List
from langgraph.graph import StateGraph, END
from langchain_openai import ChatOpenAI
from state import ContentState
import json
import matplotlib.pyplot as plt
from io import BytesIO
import base64
from dotenv import load_dotenv
from ddgs import DDGS
import wikipediaapi
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet

load_dotenv()
llm = ChatOpenAI(model="gpt-4", temperature=0.7)

def query_analyzer(state: ContentState) -> ContentState:
    print("Executing: Query Analyzer")
    prompt = f"Analyze this query and determine content type (presentation/document/webpage): {state.query}"
    response = llm.invoke(prompt)
    
    if "presentation" in response.content.lower() or "slides" in response.content.lower():
        state.content_type = "presentation"
    elif "document" in response.content.lower() or "report" in response.content.lower():
        state.content_type = "document"
    else:
        state.content_type = "webpage"
    
    print(f"Content type determined: {state.content_type}")
    return state

def research_agent(state: ContentState) -> ContentState:
    print("Executing: Research Agent")
    results = []
    
    try:
        # DuckDuckGo web search
        print("Searching DuckDuckGo...")
        ddgs = DDGS()
        web_results = ddgs.text(state.query, max_results=3)
        
        for result in web_results:
            results.append({
                'title': result.get('title', ''),
                'url': result.get('href', ''),
                'snippet': result.get('body', '')[:200] + '...' if len(result.get('body', '')) > 200 else result.get('body', '')
            })
        
        print(f"Found {len(results)} DuckDuckGo results")
        
    except Exception as e:
        print(f"DuckDuckGo search failed: {e}")
    
    try:
        # Wikipedia search
        print("Searching Wikipedia...")
        wiki = wikipediaapi.Wikipedia(
            user_agent='ContentMaster/1.0 (https://github.com/contentmaster)',
            language='en'
        )
        search_terms = state.query.split()[:2]  # Use first 2 words for better Wikipedia results
        wiki_query = ' '.join(search_terms)
        
        page = wiki.page(wiki_query)
        if page.exists():
            summary = page.summary[:300] + '...' if len(page.summary) > 300 else page.summary
            results.append({
                'title': page.title,
                'url': page.fullurl,
                'snippet': summary
            })
            print("Found Wikipedia article")
        else:
            # Try searching for related pages
            wiki_search_url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{wiki_query.replace(' ', '_')}"
            try:
                response = requests.get(wiki_search_url, timeout=5)
                if response.status_code == 200:
                    data = response.json()
                    results.append({
                        'title': data.get('title', ''),
                        'url': data.get('content_urls', {}).get('desktop', {}).get('page', ''),
                        'snippet': data.get('extract', '')[:300] + '...' if len(data.get('extract', '')) > 300 else data.get('extract', '')
                    })
                    print("Found Wikipedia summary")
            except:
                pass
                
    except Exception as e:
        print(f"Wikipedia search failed: {e}")
    
    try:
        # ArXiv API for academic papers
        print("Searching ArXiv for academic papers...")
        arxiv_url = f"http://export.arxiv.org/api/query?search_query=all:{state.query}&max_results=2"
        response = requests.get(arxiv_url, timeout=5)
        if response.status_code == 200:
            from xml.etree import ElementTree as ET
            root = ET.fromstring(response.content)
            namespace = {'atom': 'http://www.w3.org/2005/Atom'}
            
            for entry in root.findall('atom:entry', namespace)[:2]:
                title = entry.find('atom:title', namespace)
                link = entry.find('atom:id', namespace)
                summary = entry.find('atom:summary', namespace)
                
                if title is not None and link is not None:
                    results.append({
                        'title': title.text.strip(),
                        'url': link.text.strip(),
                        'snippet': summary.text.strip()[:300] + '...' if summary is not None and len(summary.text.strip()) > 300 else (summary.text.strip() if summary is not None else '')
                    })
            print(f"Found ArXiv papers")
    except Exception as e:
        print(f"ArXiv search failed: {e}")
    
    # Fallback if no results found
    if not results:
        print("Using fallback research data...")
        results = [
            {'title': f'Academic Research: {state.query}', 'url': 'https://scholar.google.com', 'snippet': f'Academic research and scholarly articles about {state.query}'},
            {'title': f'Industry Analysis: {state.query}', 'url': 'https://industry-reports.com', 'snippet': f'Industry analysis and market research for {state.query}'},
        ]
    
    state.search_results = results
    print(f"Found {len(state.search_results)} total search results")
    return state

def source_verifier(state: ContentState) -> ContentState:
    print("Executing: Source Verifier")
    verified = []
    for result in state.search_results:
        score = 0.8 if any(domain in result['url'] for domain in ['edu', 'gov', 'org']) else 0.6
        verified.append({**result, 'credibility_score': score})
    
    state.verified_sources = sorted(verified, key=lambda x: x['credibility_score'], reverse=True)
    state.quality_score = sum(s['credibility_score'] for s in state.verified_sources) / len(state.verified_sources) if state.verified_sources else 0
    
    print(f"Quality score: {state.quality_score}")
    return state

def content_planner(state: ContentState) -> ContentState:
    print("Executing: Content Planner")
    if state.content_type == "presentation":
        state.content_plan = {
            'sections': ['Title', 'Introduction', 'Main Points', 'Data/Statistics', 'Conclusion'],
            'slides': 5
        }
    elif state.content_type == "document":
        state.content_plan = {
            'sections': ['Executive Summary', 'Introduction', 'Analysis', 'Findings', 'Conclusion'],
            'pages': 3
        }
    else:
        state.content_plan = {
            'sections': ['Header', 'Overview', 'Content', 'Sources'],
            'layout': 'single_page'
        }
    
    print(f"Plan created with {len(state.content_plan['sections'])} sections")
    return state

def content_generator(state: ContentState) -> ContentState:
    print("Executing: Content Generator")
    content = {}
    sources_text = "\n".join([f"- {s['title']}: {s['snippet']}" for s in state.verified_sources[:3]])
    
    for section in state.content_plan['sections']:
        prompt = f"Generate {section} content for {state.content_type} about: {state.query}\nBased on: {sources_text}"
        response = llm.invoke(prompt)
        content[section] = response.content
    
    state.generated_content = content
    print(f"Generated content for {len(content)} sections")
    return state

def visual_creator(state: ContentState) -> ContentState:
    print("Executing: Visual Creator")
    visuals = []
    
    if state.content_type in ["presentation", "document"]:
        # Create a chart image file
        plt.figure(figsize=(10, 6))
        plt.bar(['Current Trends', 'Future Outlook', 'Key Benefits'], [30, 45, 25])
        plt.title(f"Analysis Overview: {state.query}")
        plt.ylabel('Impact Score')
        plt.xlabel('Categories')
        
        clean_query = "".join(c for c in state.query if c.isalnum() or c in (' ', '-', '_')).replace(' ', '_')
        chart_filename = f"chart_{clean_query}.png"
        plt.savefig(chart_filename, dpi=300, bbox_inches='tight')
        
        # Also save to buffer for embedding
        buffer = BytesIO()
        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
        buffer.seek(0)
        img_data = base64.b64encode(buffer.getvalue()).decode()
        plt.close()
        
        visuals.append({
            'type': 'chart',
            'filename': chart_filename,
            'data': img_data,
            'caption': f'Analysis overview for {state.query}'
        })
        
        print(f"Created chart file: {chart_filename}")
    
    state.visuals = visuals
    print(f"Created {len(visuals)} visuals")
    return state

def template_selector(state: ContentState) -> ContentState:
    print("Executing: Template Selector")
    templates = {
        "presentation": "modern_slides",
        "document": "professional_report",
        "webpage": "clean_web"
    }
    state.template = templates.get(state.content_type, "default")
    print(f"Selected template: {state.template}")
    return state

def content_assembler(state: ContentState) -> ContentState:
    print("Executing: Content Assembler")
    output = {
        'type': state.content_type,
        'template': state.template,
        'content': state.generated_content,
        'visuals': state.visuals,
        'sources': [s['url'] for s in state.verified_sources],
        'metadata': {
            'query': state.query,
            'quality_score': state.quality_score
        }
    }
    
    # Create actual files based on content type
    files_created = []
    
    if state.content_type == "presentation":
        filename = create_presentation_file(state)
        files_created.append(filename)
        print(f"Created PowerPoint presentation: {filename}")
        
    elif state.content_type == "document":
        filename = create_document_file(state)
        files_created.append(filename)
        print(f"Created PDF document: {filename}")
        
    elif state.content_type == "webpage":
        filename = create_webpage_file(state)
        files_created.append(filename)
        print(f"Created HTML webpage: {filename}")
    
    output['files_created'] = files_created
    state.final_output = output
    print(f"Content assembly completed! Created {len(files_created)} files")
    return state

def create_presentation_file(state: ContentState) -> str:
    """Create actual PowerPoint presentation file"""
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = state.query.title()
    subtitle.text = f"Research-backed presentation\nSources: {len(state.verified_sources)}"
    
    # Content slides
    for section, content in state.generated_content.items():
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = section
        tf = body_shape.text_frame
        tf.text = content[:500] + "..." if len(content) > 500 else content
    
    # Add chart if available
    if state.visuals:
        chart_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(chart_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Data Analysis"
        
        chart_filename = state.visuals[0].get('filename')
        if chart_filename and os.path.exists(chart_filename):
            slide.shapes.add_picture(chart_filename, Inches(1), Inches(2), Inches(8), Inches(4))
    
    # Sources slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = "Sources"
    tf = body_shape.text_frame
    for i, source in enumerate(state.verified_sources[:5]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {source['title']}: {source['url']}"
    
    clean_query = "".join(c for c in state.query if c.isalnum() or c in (' ', '-', '_')).replace(' ', '_')
    filename = f"{clean_query}_presentation.pptx"
    prs.save(filename)
    return filename

def create_document_file(state: ContentState) -> str:
    """Create actual PDF document file"""
    clean_query = "".join(c for c in state.query if c.isalnum() or c in (' ', '-', '_')).replace(' ', '_')
    filename = f"{clean_query}_document.pdf"
    doc = SimpleDocTemplate(filename, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    title = Paragraph(state.query.title(), styles['Title'])
    story.append(title)
    story.append(Spacer(1, 12))
    
    # Content sections
    for section, content in state.generated_content.items():
        section_title = Paragraph(section, styles['Heading2'])
        story.append(section_title)
        
        content_para = Paragraph(content, styles['Normal'])
        story.append(content_para)
        story.append(Spacer(1, 12))
    
    # Add chart if available
    if state.visuals:
        chart_filename = state.visuals[0].get('filename')
        if chart_filename and os.path.exists(chart_filename):
            story.append(Paragraph("Data Analysis", styles['Heading2']))
            img = Image(chart_filename, width=400, height=240)
            story.append(img)
            story.append(Spacer(1, 12))
    
    # Sources
    story.append(Paragraph("Sources", styles['Heading2']))
    for source in state.verified_sources:
        source_para = Paragraph(f"• {source['title']}: {source['url']}", styles['Normal'])
        story.append(source_para)
    
    doc.build(story)
    return filename

def create_webpage_file(state: ContentState) -> str:
    """Create actual HTML webpage file"""
    clean_query = "".join(c for c in state.query if c.isalnum() or c in (' ', '-', '_')).replace(' ', '_')
    filename = f"{clean_query}_webpage.html"
    
    html_content = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>{state.query.title()}</title>
        <style>
            body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
            h1 {{ color: #333; border-bottom: 2px solid #007acc; }}
            h2 {{ color: #007acc; }}
            .section {{ margin-bottom: 30px; }}
            .sources {{ background: #f5f5f5; padding: 20px; }}
            img {{ max-width: 100%; height: auto; }}
        </style>
    </head>
    <body>
        <h1>{state.query.title()}</h1>
        <p><em>Research-backed content from {len(state.verified_sources)} sources</em></p>
    """
    
    # Content sections
    for section, content in state.generated_content.items():
        html_content += f"""
        <div class="section">
            <h2>{section}</h2>
            <p>{content}</p>
        </div>
        """
    
    # Add chart if available
    if state.visuals:
        chart_filename = state.visuals[0].get('filename')
        if chart_filename and os.path.exists(chart_filename):
            html_content += f"""
            <div class="section">
                <h2>Data Analysis</h2>
                <img src="{chart_filename}" alt="Analysis Chart">
            </div>
            """
    
    # Sources
    html_content += """
        <div class="sources">
            <h2>Sources</h2>
            <ul>
    """
    for source in state.verified_sources:
        html_content += f'<li><a href="{source["url"]}">{source["title"]}</a></li>'
    
    html_content += """
            </ul>
        </div>
    </body>
    </html>
    """
    
    with open(filename, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    return filename

def should_retry_research(state: ContentState) -> str:
    print(f"Research check: results={len(state.search_results)}")
    return "proceed"

def should_retry_verification(state: ContentState) -> str:
    print(f"Verification check: quality_score={state.quality_score}")
    return "proceed"

def needs_visuals(state: ContentState) -> str:
    decision = "with_visuals" if state.content_type in ["presentation", "document"] else "no_visuals"
    print(f"Visuals check: content_type={state.content_type}, decision={decision}")
    return decision

def create_workflow():
    workflow = StateGraph(ContentState)
    
    workflow.add_node("query_analyzer", query_analyzer)
    workflow.add_node("research_agent", research_agent)
    workflow.add_node("source_verifier", source_verifier)
    workflow.add_node("content_planner", content_planner)
    workflow.add_node("content_generator", content_generator)
    workflow.add_node("visual_creator", visual_creator)
    workflow.add_node("template_selector", template_selector)
    workflow.add_node("content_assembler", content_assembler)
    
    workflow.set_entry_point("query_analyzer")
    
    workflow.add_edge("query_analyzer", "research_agent")
    workflow.add_conditional_edges(
        "research_agent", 
        should_retry_research,
        {"retry": "research_agent", "proceed": "source_verifier"}
    )
    workflow.add_conditional_edges(
        "source_verifier",
        should_retry_verification,
        {"retry": "research_agent", "proceed": "content_planner"}
    )
    workflow.add_edge("content_planner", "content_generator")
    workflow.add_conditional_edges(
        "content_generator",
        needs_visuals,
        {"with_visuals": "visual_creator", "no_visuals": "template_selector"}
    )
    workflow.add_edge("visual_creator", "template_selector")
    workflow.add_edge("template_selector", "content_assembler")
    workflow.add_edge("content_assembler", END)
    
    return workflow.compile()

def run_content_master(query: str):
    app = create_workflow()
    initial_state = ContentState(query=query)
    config = {"recursion_limit": 50}
    result = app.invoke(initial_state, config=config)
    return result 